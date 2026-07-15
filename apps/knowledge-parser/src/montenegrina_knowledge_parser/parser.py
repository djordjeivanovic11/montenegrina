from __future__ import annotations

import io
import re
from typing import Any

from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from markdown import markdown
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from .sanitization import sanitize_metadata, sanitize_text

ARTICLE_PATTERN = re.compile(r"(?i)^(?:član|clan|article)\s+(\d+[a-z]?)\b")


def detect_article_number(heading: str | None, content: str) -> str | None:
    for candidate in filter(None, [heading, content.split("\n", 1)[0] if content else None]):
        match = ARTICLE_PATTERN.search(candidate.strip())
        if match:
            return match.group(1)
    return None


def section(
    *,
    heading: str | None,
    content: str,
    level: int = 0,
    page_start: int | None = None,
    page_end: int | None = None,
    parent_index: int | None = None,
    is_table: bool = False,
    metadata: dict[str, Any] | None = None,
) -> dict[str, Any]:
    clean_heading = sanitize_text(heading) if heading is not None else None
    clean_content = sanitize_text(content).strip()
    return {
        "heading": clean_heading,
        "level": level,
        "pageStart": page_start,
        "pageEnd": page_end,
        "articleNumber": detect_article_number(clean_heading, clean_content),
        "content": clean_content,
        "parentIndex": parent_index,
        "isTable": is_table,
        "metadata": sanitize_metadata(metadata or {}),
    }


def parse_pdf(data: bytes) -> list[dict[str, Any]]:
    reader = PdfReader(io.BytesIO(data))
    sections: list[dict[str, Any]] = []
    for page_number, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if not text:
            continue
        blocks = [block.strip() for block in re.split(r"\n{2,}", text) if block.strip()]
        for block in blocks:
            heading = block.split("\n", 1)[0][:120] if len(block) > 80 else None
            sections.append(
                section(
                    heading=heading,
                    content=block,
                    level=1,
                    page_start=page_number,
                    page_end=page_number,
                )
            )
    return sections


def parse_docx(data: bytes) -> list[dict[str, Any]]:
    document = DocxDocument(io.BytesIO(data))
    sections: list[dict[str, Any]] = []
    current_heading: str | None = None
    current_parts: list[str] = []
    parent_index: int | None = None

    def flush() -> None:
        nonlocal current_heading, current_parts, parent_index
        body = "\n".join(current_parts).strip()
        if body:
            sections.append(
                section(
                    heading=current_heading,
                    content=body,
                    level=2 if current_heading else 0,
                    parent_index=parent_index,
                )
            )
        current_parts = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = (paragraph.style.name or "").lower()
        if "heading" in style:
            flush()
            current_heading = text
            parent_index = len(sections) - 1 if sections else None
            continue
        current_parts.append(text)
    flush()
    return sections or [section(heading=None, content="\n".join(p.text for p in document.paragraphs if p.text))]


def parse_pptx(data: bytes) -> list[dict[str, Any]]:
    presentation = Presentation(io.BytesIO(data))
    sections: list[dict[str, Any]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts = []
        title = None
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if not text:
                continue
            if title is None:
                title = text
            else:
                parts.append(text)
        body = "\n".join(parts).strip()
        if title or body:
            sections.append(
                section(
                    heading=title,
                    content=body or title or "",
                    level=1,
                    page_start=index,
                    page_end=index,
                )
            )
    return sections


def parse_xlsx(data: bytes) -> list[dict[str, Any]]:
    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sections: list[dict[str, Any]] = []
    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if values:
                rows.append(" | ".join(values))
        if rows:
            sections.append(
                section(
                    heading=sheet.title,
                    content="\n".join(rows),
                    level=1,
                    is_table=True,
                )
            )
    return sections


def parse_text(data: bytes, media_type: str) -> list[dict[str, Any]]:
    text = data.decode("utf-8")
    if media_type in {"text/markdown", "text/x-markdown"}:
        text = BeautifulSoup(markdown(text), "html.parser").get_text("\n")
    if media_type in {"text/html", "application/xhtml+xml"}:
        text = BeautifulSoup(text, "html.parser").get_text("\n")
    blocks = [block.strip() for block in re.split(r"\n{2,}", text) if block.strip()]
    sections: list[dict[str, Any]] = []
    for block in blocks:
        lines = block.split("\n")
        heading = lines[0] if len(lines) > 1 and len(lines[0]) < 120 else None
        content = block if heading is None else "\n".join(lines[1:]).strip() or block
        sections.append(section(heading=heading, content=content, level=1 if heading else 0))
    return sections


def parse_document(data: bytes, media_type: str) -> dict[str, Any]:
    media_type = media_type.split(";")[0].strip().lower()
    if media_type == "application/pdf":
        sections = parse_pdf(data)
    elif media_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        sections = parse_docx(data)
    elif media_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        sections = parse_pptx(data)
    elif media_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        sections = parse_xlsx(data)
    elif media_type in {"text/plain", "text/markdown", "text/x-markdown", "text/html", "application/xhtml+xml"}:
        sections = parse_text(data, media_type)
    else:
        raise ValueError(f"UNSUPPORTED_MEDIA_TYPE:{media_type}")

    if not sections:
        raise ValueError("DOCUMENT_EMPTY")

    page_count = max(
        [section.get("pageEnd") or section.get("pageStart") or 0 for section in sections] or [0]
    )
    extracted_text = "\n\n".join(section["content"] for section in sections if section["content"])
    return {
        "parserVersion": "montenegrina-knowledge-parser-0.2.0",
        "pageCount": page_count or None,
        "extractedText": extracted_text,
        "sections": sections,
    }
